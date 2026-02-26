#!/usr/bin/env python3
"""
Build ProteinProcessIO presentation (.pptx) from PRESENTATION.md and GP15 calibration note.

Usage:
  pip install python-pptx
  python scripts/build_presentation.py

Output: ProteinProcessIO_Presentation.pptx (project root)
"""
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    raise SystemExit(f"Install python-pptx: pip install python-pptx\n{e}") from e

# Slide dimensions (default 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_TOP = Inches(1.2)
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN


def set_shape_bullets(shape, level0_size=18, level1_size=14):
    """Set bullet paragraph font sizes for a text frame."""
    for para in shape.text_frame.paragraphs:
        para.font.size = Pt(level0_size)
        para.space_after = Pt(6)
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(level0_size)


def add_title_slide(prs, title, subtitle_lines):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Title
    box = slide.shapes.add_textbox(MARGIN, Inches(2.2), CONTENT_WIDTH, Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    box2 = slide.shapes.add_textbox(MARGIN, Inches(3.6), CONTENT_WIDTH, Inches(1.5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(subtitle_lines):
        p2 = tf2.paragraphs[i] if i < len(tf2.paragraphs) else tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(18)
        p2.font.italic = True
        p2.alignment = PP_ALIGN.CENTER
        p2.space_after = Pt(4)
    return slide


def add_section_slide(prs, title, bullets, sub_bullets=None):
    """One main title + bullet list. sub_bullets: list of lists, one per main bullet."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # Section title
    box = slide.shapes.add_textbox(MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(0.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    # Content
    box2 = slide.shapes.add_textbox(MARGIN, CONTENT_TOP, CONTENT_WIDTH, Inches(6.2))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    sub_bullets = sub_bullets or []
    for i, bullet in enumerate(bullets):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        para.text = bullet
        para.font.size = Pt(16)
        para.space_after = Pt(8)
        para.level = 0
        if i < len(sub_bullets) and sub_bullets[i]:
            for sub in sub_bullets[i]:
                p2 = tf2.add_paragraph()
                p2.text = sub
                p2.font.size = Pt(14)
                p2.space_after = Pt(4)
                p2.level = 1
    return slide


def add_bullet_slide(prs, title, bullets, level1=None):
    """Simple title + bullets. level1 = list of (index, sub-bullets)."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    box = slide.shapes.add_textbox(MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(0.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    box2 = slide.shapes.add_textbox(MARGIN, CONTENT_TOP, CONTENT_WIDTH, Inches(6.2))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    level1 = level1 or {}
    for i, bullet in enumerate(bullets):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        para.text = bullet
        para.font.size = Pt(15)
        para.space_after = Pt(6)
        para.level = 0
        if i in level1:
            for sub in level1[i]:
                p2 = tf2.add_paragraph()
                p2.text = "  • " + sub
                p2.font.size = Pt(13)
                p2.space_after = Pt(2)
                p2.level = 1
    return slide


def add_table_slide(prs, title, headers, rows, col_widths=None):
    """Add a slide with a title and a table."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    box = slide.shapes.add_textbox(MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(0.6))
    tf = box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    ncols = len(headers)
    nrows = len(rows) + 1
    left = MARGIN
    top = Inches(1.3)
    width = CONTENT_WIDTH
    height = min(Inches(0.35 * nrows), Inches(4))
    if col_widths is None:
        col_widths = [width / ncols] * ncols
    table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(h)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            if c < ncols:
                table.cell(r + 1, c).text = str(val)[:80]
                table.cell(r + 1, c).text_frame.paragraphs[0].font.size = Pt(10)
    return slide


def main():
    out_path = Path(__file__).resolve().parent.parent / "ProteinProcessIO_Presentation.pptx"
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # --- Slide 1: Title ---
    add_title_slide(
        prs,
        "ProteinProcessIO",
        [
            "A GPU-Accelerated Digital Twin for Multi-Stage Dry Air Classification of Plant Proteins",
            "High-Fidelity Multiphysics Simulation | NVIDIA Warp GPU | Interactive Desktop App",
            "From Thermal Pretreatment to Protein Separation — End-to-End Process Design",
        ],
    )

    # --- Slide 2: Introduction ---
    add_section_slide(
        prs,
        "Introduction — What is ProteinProcessIO?",
        [
            "Physics-based digital twin for the complete dry fractionation pipeline: RF Thermal Pretreatment, Pin Milling, Air Classification.",
            "Interactive desktop application with real-time 3D visualization.",
            "Enables design, simulate, optimize, and validate industrial classification systems in software before physical prototypes.",
            "Why it matters: Plant-based protein demand is growing; dry fractionation is the most sustainable extraction (no water, no chemicals). ProteinProcessIO replaces expensive trial-and-error with computer-aided engineering.",
        ],
    )

    # --- Slide 3: The Problem ---
    add_section_slide(
        prs,
        "The Problem — Challenges in Dry Protein Fractionation",
        [
            "Complex, coupled physics — particle separation depends on aerodynamics, PSD, material properties, and geometry interacting simultaneously.",
            "No integrated simulation tools — engineers rely on spreadsheets, correlations, or single-component CFD that miss cross-stage effects.",
            "Expensive physical prototyping — each pilot trial costs thousands; iterating is slow.",
            "Multi-stage coupling — pretreatment affects milling, which affects PSD, which determines classification. Optimizing one stage in isolation is suboptimal.",
            "Competing objectives — protein purity vs. yield; thermal treatment for flavor vs. protein denaturation risk.",
            "The gap: No software combines parametric geometry + multiphysics + GPU + multi-objective optimization for food powder classification.",
        ],
    )

    # --- Slide 4: Objective, Research Questions, Hypothesis ---
    add_section_slide(
        prs,
        "Objective, Research Questions & Hypothesis",
        [
            "Objective: Develop a GPU-accelerated, physics-based digital twin for computer-aided design and optimization of multi-stage dry air classification — from thermal pretreatment through milling to protein–starch separation.",
            "Research questions: (1) Can coupled multiphysics (RF + particle dynamics + fluid) predict real equipment when calibrated with sensor data? (2) Can GPU (NVIDIA Warp) make 3D simulation fast enough for interactive exploration? (3) Can multi-objective optimization (Derringer–Suich) identify optimal conditions for yield, purity, energy, and quality?",
            "Hypothesis: A physics-based digital twin calibrated against PLC sensor data will predict classification performance within engineering accuracy (<15% error), with GPU enabling interactive runs (<5 min per full simulation).",
        ],
    )

    # --- Slide 5: Methods Overview ---
    add_section_slide(
        prs,
        "Methods Overview — Digital Twin Methodology",
        [
            "Parametric 3D geometry (40+ parts) → Multiphysics simulation (GPU-accelerated) → Calibration with PLC sensor data.",
            "Outputs feed: Interactive desktop GUI ← Multi-objective optimization ← Validated digital twin.",
            "Why physics-based over ML: Lower data requirement; reliable extrapolation; interpretable; supports new materials and scale-up; design exploration; regulatory trust.",
        ],
    )

    # --- Slide 6: Multiphysics — Three Stages ---
    add_section_slide(
        prs,
        "Methods — Three-Stage Coupled Pipeline",
        [
            "Stage 1 — RF Pretreatment (GP-15): Laplace (E-field), heat equation (T), diffusion (M), dielectric properties, PLC logic, belt advection. 9-step coupling per timestep.",
            "Stage 2 — Pin Milling: Impact detection, breakage model (selection + breakage function), screen passage, population balance.",
            "Stage 3 — Air Classification: Zone-based Lagrangian tracking — Venturi eductor, zigzag classifier, wheel classifier (300–5000 g), 3-stage cyclones, bag filter. Drag: Stokes, Schiller–Naumann, Haider–Levenspiel.",
        ],
    )

    # --- Slide 7: Calibration with Real Hardware ---
    add_section_slide(
        prs,
        "Methods — Calibration with Real Hardware (GP-15)",
        [
            "Strategy: Match PLC sensors (temperature, RF current, belt speed, moisture) and temp strips to simulation under same operating conditions; tune oscillator efficiency, contact resistance, convection.",
            "GP-15 Run#2 validation: 90 kg yellow pea, 35 mm bed, 0.2 m/min, electrode gap 75→94 mm. Outfeed 68–82°C, moisture loss 1.3 pp, LOX inactivation 24.2 min above 65°C, specific energy 0.042 kWh/kg.",
            "Calibration store: Empirical curves and tuned parameters saved for reuse across operating conditions.",
        ],
    )

    # --- Slide 8: Calibration Innovations (from Technical Note) ---
    add_table_slide(
        prs,
        "Calibration Innovations — GP-15 Digital Twin (Technical Note)",
        ["Result", "Detail"],
        [
            ("Four-parameter calibration", "k_c, k_evap, r_gap, k_disp via differential evolution + L-BFGS-B on full PLC time-series"),
            ("Electrode gap prediction", "Run#2 blind: 93.0 mm sim vs 94.1 mm PLC — 1.2% error; validates EM + control chain"),
            ("Energy partition", "99.2% sensible heating, ~0.8% latent; seed coat suppresses evaporation at sub-100°C"),
            ("MRH edge regime", "Anode current ~1.70 A near threshold; proportional gap drift correctly predicted"),
            ("Sensor-comparable metric", "75th percentile outfeed T for PLC/strip comparison; P75 vs mean vs max"),
            ("GPU performance", "57,600 cells; ~3 s/eval; calibration 926 evals in ~46 min; 60× speedup with array reuse"),
        ],
    )

    # --- Slide 9: Technology Stack ---
    add_section_slide(
        prs,
        "Methods — Python Technology Stack",
        [
            "GUI: PySide6 (Qt6) + PyVistaQt — 3D viewport, mode switching (Classification / Pretreatment / Milling), cinematic camera, live KPI dashboards.",
            "Visualization: PyVista (VTK) — 40+ parametric meshes, real-time particles, physics-driven animations.",
            "Simulation: NumPy, SciPy, NVIDIA Warp — coupled solvers, Lagrangian tracking, zone-based classification, SPH air flow.",
            "GPU: NVIDIA Warp (JIT CUDA) — 10+ kernels, persistent GPU memory, batched launch, single sync.",
            "Data: Pandas, Matplotlib, PyYAML — CSV/JSON/VTK export, publication plots. Packaging: PyInstaller for standalone executable.",
        ],
    )

    # --- Slide 10: NVIDIA Warp GPU ---
    add_section_slide(
        prs,
        "Methods — NVIDIA Warp GPU Acceleration",
        [
            "Warp: Python framework; @wp.kernel functions JIT-compiled to CUDA; vec3, mat33, Mesh, HashGrid; automatic CPU/CUDA fallback.",
            "Kernels per timestep (batched, single sync): RF field, thermal, moisture, dielectric, transport, particle (drag, collisions, zone tracking).",
            "Strategy: Persistent GPU arrays, batched launch, one wp.synchronize() per step, kernel cache. Advantage vs raw CUDA: pure Python syntax, fast iteration, NumPy/SciPy interop, maintainable by non-GPU specialists.",
        ],
    )

    # --- Slide 11: Multi-Objective Optimization ---
    add_section_slide(
        prs,
        "Methods — Multi-Objective Optimization (Derringer–Suich)",
        [
            "Pretreatment: 5 objectives — thermal treatment (65–82°C), LOX kill (≥4 min above 65°C), protein preservation (peak <71°C vicilin), moisture retention (<2 pp loss), energy (<0.15 kWh/kg). Overall desirability = geometric mean; scaled 0–10.",
            "Material-specific profiles: Yellow pea (7S <71°C), faba bean (11S <90°C), red lentil (<75°C).",
            "Classification: Design variables — wheel RPM, cyclone diameters, air flow. Objectives — maximize yield and purity, minimize energy; constraints on efficiency and starch purity.",
        ],
    )

    # --- Slide 12: Environmental & Economics ---
    add_bullet_slide(
        prs,
        "Methods — Integrated Environmental & Economics",
        [
            "End-to-end: Raw seeds → RF Pretreatment → Pin Milling → Air Classification. Energy per stage tracked; zero water, zero chemical waste.",
            "Environmental: Energy per kg, carbon from energy source; economic: capex from parametric geometry, energy cost per kg, yield (protein vs starch), throughput for capacity planning.",
            "Dry vs wet: Dry ~0.05–0.10 kWh/kg, 0 L water/kg, 50–65% purity; wet ~0.3–0.5 kWh/kg, 5–10 L water/kg, 80–90% purity. Digital twin enables per-stage accounting and Pareto yield–purity curves.",
        ],
    )

    # --- Slide 13: Desktop Application ---
    add_section_slide(
        prs,
        "Methods — Desktop Application (PyInstaller)",
        [
            "Standalone executable; three-mode interface (Classification, Pretreatment, Milling); real-time 3D viewport (~60 FPS).",
            "Control panel: Material, feed rate, wheel RPM, air flow; live KPIs (efficiency, yield, purity, power); export CSV/VTK.",
            "Physics-driven animations (blower, damper, wheel, belt); cinematic camera (orbit, showcase, flythrough); assembly configurator; live PSD and efficiency curves.",
        ],
    )

    # --- Slide 14: Results — Validation ---
    add_table_slide(
        prs,
        "Results — Pretreatment Validation (GP-15 Run#2)",
        ["Metric", "Target", "Achieved", "Status"],
        [
            ("Outlet temperature", "65–82°C", "68–82°C (PLC + strips)", "PASS"),
            ("LOX inactivation", "≥4 min above 65°C", "24.2 min", "PASS"),
            ("Protein preservation", "Max T < 84°C (legumin)", "Legumin intact (11S)", "PASS"),
            ("Moisture retention", "<2 pp loss", "1.3 pp (11.8%→10.5%)", "PASS"),
            ("Specific energy", "<0.15 kWh/kg", "0.042 kWh/kg", "PASS"),
        ],
    )

    # --- Slide 15: Results — Classification & Geometry ---
    add_bullet_slide(
        prs,
        "Results — Classification & Parametric Geometry",
        [
            "Classification (simulated): Cut size ~25 µm @ 3000 rpm; protein yield ~30%, purity ~55%; starch purity ~70%; separation efficiency ~45%.",
            "Parametric geometry: 40+ industrial components, zero hardcoded dimensions; pilot (50 kg/h) to production (2000+ kg/h); dual mesh/SDF for rendering and GPU physics.",
            "GPU: CUDA/CPU auto-detection, persistent memory, batched kernels, <1 ms CPU overhead per step.",
        ],
    )

    # --- Slide 16: Key Innovations ---
    add_section_slide(
        prs,
        "Key Innovations & Breakthroughs",
        [
            "First physics-based digital twin for full dry fractionation (pretreatment + milling + classification) in one coupled simulation.",
            "Whole-seed RF thermal conditioning model: series-capacitor voltage division, Lagrangian tracers; whole seeds show minimal drying (1–2% vs 3–4% for flour).",
            "GPU-accelerated multiphysics in pure Python: 10+ Warp kernels (EM, thermal, moisture, dielectric, particles, advection); near-native CUDA with NumPy/SciPy interop.",
            "Parametric geometry framework: 40+ components from design equations; dual mesh + SDF; automatic scaling.",
            "Multi-objective Desirability (Derringer–Suich); material-specific profiles; interactive desktop with real-time 3D and physics-driven animations.",
            "Calibrated against real GP-15 with PLC data; calibration store for reuse. Comprehensive technical documentation (geometry, multiphysics, Warp, pretreatment guide).",
        ],
    )

    # --- Slide 17: Summary & Next Steps ---
    add_table_slide(
        prs,
        "Summary & Next Steps",
        ["Module", "Status"],
        [
            ("RF Pretreatment (GP-15)", "Functional — calibrating against Run#2"),
            ("Hammer Mill (Milling)", "In development — geometry and physics in place"),
            ("Air Classification", "Functional — zone-based physics"),
            ("GUI Application", "Functional — 3-mode interface, 3D viewport"),
            ("Optimization", "Framework implemented — expanding to full pipeline"),
        ],
    )

    add_bullet_slide(
        prs,
        "Next Steps",
        [
            "Complete milling module — breakage mechanics and screen classification.",
            "Full pipeline coupling — end-to-end from seeds to separated protein.",
            "Extended calibration — more experimental runs for validation.",
            "Web documentation portal — installation, user manual, API.",
            "Industrial partner trials — validate at production scale.",
            "Additional materials — oat, chickpea, lentil varieties.",
        ],
    )

    # --- Slide 18: Thank You ---
    add_title_slide(
        prs,
        "ProteinProcessIO",
        [
            "Accelerating the future of sustainable plant protein processing through physics-based digital twin technology.",
            "Questions?",
        ],
    )

    prs.save(str(out_path))
    print(f"Saved: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
